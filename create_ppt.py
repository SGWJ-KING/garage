from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1B, 0x2A, 0x3A)   # 深蓝黑背景
ACCENT     = RGBColor(0x2E, 0x86, 0xC1)   # 蓝色主色调
ACCENT2    = RGBColor(0xE7, 0x4C, 0x3C)   # 红色强调
ACCENT3    = RGBColor(0x27, 0xAE, 0x60)   # 绿色通过
ACCENT4    = RGBColor(0xF3, 0x9C, 0x12)   # 橙色待定
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_TEXT   = RGBColor(0x2C, 0x3E, 0x50)
MID_GRAY   = RGBColor(0x95, 0xA5, 0xA6)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bottom_bar(slide):
    left, top, w, h = Inches(0), Inches(7.1), Inches(13.333), Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_top_accent_line(slide):
    left, top, w, h = Inches(0), Inches(0), Inches(13.333), Inches(0.06)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT


def add_title_shape(slide, title_text, subtitle_text=None):
    """Main-content title bar at the top of a content slide."""
    left, top, w, h = Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.3)

    if subtitle_text:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(12), Inches(0.5))
        stf = txBox.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(14)
        sp.font.color.rgb = MID_GRAY
        sp.font.italic = True


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 ── 封面
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)

# big title
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "简历筛选与评估模板"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# subtitle
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Resume Screening & Evaluation Framework"
p.font.size = Pt(22)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# decorative line
left, top, w, h = Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.04)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# info line
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.0))
tf = txBox.text_frame
for txt in ["适用岗位: _______________", "日期: _______________", "筛选负责人: _______________"]:
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 ── 筛选流程概览
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "简历筛选流程概览", "Resume Screening Process Overview")
add_slide_number(slide, 2)

steps = [
    ("01", "简历收集", "通过招聘渠道\n收集候选人简历"),
    ("02", "初筛过滤", "根据硬性条件\n快速过滤不合格简历"),
    ("03", "细评打分", "对通过初筛的简历\n进行详细评估打分"),
    ("04", "对比排序", "横向对比候选人\n按得分高低排序"),
    ("05", "面试邀约", "确定面试名单\n发出面试邀请"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.7 + i * 2.45)
    # circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.2), Inches(1.2), Inches(1.2))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # arrow (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.35), Inches(2.6), Inches(0.95), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

    # title below
    txBox = slide.shapes.add_textbox(x - Inches(0.1), Inches(3.7), Inches(1.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # desc
    txBox = slide.shapes.add_textbox(x - Inches(0.1), Inches(4.2), Inches(1.5), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 ── 基本信息筛选标准
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "基本信息筛选标准", "硬性条件 — 任一项不满足即淘汰")
add_slide_number(slide, 3)

criteria_data = [
    ("学历要求", "□ 本科及以上  □ 硕士及以上  □ 博士\n专业方向: _______________", "必须满足"),
    ("工作经验", "□ 应届生  □ 1-3年  □ 3-5年  □ 5-10年  □ 10年以上\n行业背景: _______________", "必须满足"),
    ("专业技能", "核心技术: _______________\n工具/软件: _______________\n证书要求: _______________", "至少满足3项"),
    ("语言能力", "□ CET-4  □ CET-6  □ 专业八级\n□ 其他语种: _______________\n□ 无硬性要求", "视岗位而定"),
    ("薪资期望", "候选人期望: _______________\n岗位预算范围: _______________\n□ 超出预算则淘汰", "预算内"),
]

for i, (item, content, tag) in enumerate(criteria_data):
    y = Inches(2.0 + i * 1.05)
    # label
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(2.0), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # content
    txBox = slide.shapes.add_textbox(Inches(2.7), y - Inches(0.05), Inches(7.5), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.line_spacing = Pt(18)

    # tag
    tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), y + Inches(0.05), Inches(2.3), Inches(0.35))
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = ACCENT2 if "必须" in tag else ACCENT4
    tag_shape.line.fill.background()
    tf = tag_shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 ── 核心能力评估
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "核心能力评估维度", "评分: 1(差) → 5(优秀)")
add_slide_number(slide, 4)

abilities = [
    ("专业能力", "专业知识深度与广度、行业理解", "40%"),
    ("项目经验", "相关项目数量、复杂度、角色贡献", "25%"),
    ("学习能力", "新技术掌握速度、自驱力", "15%"),
    ("沟通协作", "团队合作、跨部门协调能力", "10%"),
    ("稳定性/潜力", "职业规划清晰度、成长潜力", "10%"),
]

# table header
header_y = Inches(2.0)
for j, (hdr, w) in enumerate([("评估维度", 3.0), ("评估要点", 6.5), ("权重", 1.5), ("评分(1-5)", 1.5)]):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6 + sum([3.0, 6.5, 1.5, 1.5][:j]) if j > 0 else 0.6),
        header_y,
        Inches(w),
        Inches(0.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = DARK_BG
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = hdr
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

col_widths = [3.0, 6.5, 1.5, 1.5]
col_starts = [0.6, 3.6, 10.1, 11.6]

for i, (dim, points, weight) in enumerate(abilities):
    y = Inches(2.5 + i * 0.85)
    row_data = [dim, points, weight, ""]
    for j, (txt, width, x) in enumerate(zip(row_data, col_widths, col_starts)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), y, Inches(width), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x23, 0x34, 0x45) if i % 2 == 0 else RGBColor(0x2C, 0x3E, 0x50)
        shape.line.color.rgb = DARK_BG
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(13)
        p.font.bold = True if j == 0 else False
        p.font.color.rgb = WHITE if j != 3 else MID_GRAY
        p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
        if j <= 1:
            tf.margin_left = Inches(0.15)

# total row
total_y = Inches(2.5 + 5 * 0.85)
for j, (txt, w, x) in enumerate(zip(["总分", "", "100%", "_____ / 25"], col_widths, col_starts)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), total_y, Inches(w), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = DARK_BG
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 ── 加分项与减分项
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "加分项与减分项", "Bonus / Penalty Items")
add_slide_number(slide, 5)

# left ── 加分项
left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5))
left_shape.fill.solid()
left_shape.fill.fore_color.rgb = RGBColor(0x20, 0x3A, 0x2C)  # dark green tint
left_shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "✅  加分项 (+1 ~ +3 分)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT3

bonus_items = [
    "知名大厂 / 行业头部企业工作经历",
    "拥有相关专利或发表过技术文章",
    "有创业经历或参与开源项目",
    "持有含金量高的专业认证 (PMP/CPA等)",
    "海外留学或跨国项目工作经验",
    "面试过程中展现出极强的主动性",
    "推荐信 / 行业内良好口碑",
]

for i, item in enumerate(bonus_items):
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.7 + i * 0.52), Inches(5.0), Inches(0.5))
    tf = txBox2.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {item}"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY

# right ── 减分项
right_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.5))
right_shape.fill.solid()
right_shape.fill.fore_color.rgb = RGBColor(0x3D, 0x1E, 0x1E)
right_shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "❌  减分项 (-1 ~ -3 分)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ACCENT2

penalty_items = [
    "频繁跳槽 (平均每份工作不到1年)",
    "简历信息与实际不符 / 造假嫌疑",
    "职业空窗期过长且无合理解释",
    "履历中存在明显倒退 (职级/薪资)",
    "面试态度消极、不配合或迟到",
    "缺乏基本的职业素养和礼仪",
    "薪资期望严重偏离市场合理水平",
    "过往项目成果模糊、无法量化",
]

for i, item in enumerate(penalty_items):
    txBox2 = slide.shapes.add_textbox(Inches(7.4), Inches(2.7 + i * 0.5), Inches(5.0), Inches(0.5))
    tf = txBox2.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 ── 综合评分表
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "综合评分表 (模板示例)", "三位候选人横向对比")
add_slide_number(slide, 6)

eval_headers = ["评估项目", "权重", "候选人A", "候选人B", "候选人C"]
eval_col_widths = [3.0, 1.2, 2.8, 2.8, 2.8]
eval_data = [
    ("学历匹配度", "10%", "___", "___", "___"),
    ("工作经验匹配度", "25%", "___", "___", "___"),
    ("专业技能匹配度", "30%", "___", "___", "___"),
    ("项目经验相关性", "15%", "___", "___", "___"),
    ("综合素质", "10%", "___", "___", "___"),
    ("加分项合计", "5%", "___", "___", "___"),
    ("减分项合计", "5%", "___", "___", "___"),
]

# header
for j, (hdr, w) in enumerate(zip(eval_headers, eval_col_widths)):
    x = Inches(0.5 + sum(eval_col_widths[:j]))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(w), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = DARK_BG
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = hdr
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

for i, row in enumerate(eval_data):
    y = Inches(2.6 + i * 0.55)
    for j, (txt, w) in enumerate(zip(row, eval_col_widths)):
        x = Inches(0.5 + sum(eval_col_widths[:j]))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(w), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x23, 0x34, 0x45) if i % 2 == 0 else RGBColor(0x2C, 0x3E, 0x50)
        shape.line.color.rgb = DARK_BG
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(12)
        p.font.bold = (j == 0 or j == 1)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# total row
total_y2 = Inches(2.6 + 7 * 0.55)
total_row = ["综合总分", "100%", "___", "___", "___"]
for j, (txt, w) in enumerate(zip(total_row, eval_col_widths)):
    x = Inches(0.5 + sum(eval_col_widths[:j]))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, total_y2, Inches(w), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT3
    shape.line.color.rgb = DARK_BG
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# remark at bottom
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "备注: 每项按 1-5 分评分, 得分 = 评分 × 权重比例, 综合总分 = 各项得分之和"
p.font.size = Pt(11)
p.font.color.rgb = MID_GRAY

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 ── 面试决策矩阵
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)
add_title_shape(slide, "面试决策矩阵", "筛选结论与下一步行动")
add_slide_number(slide, 7)

decisions = [
    (ACCENT3, "≥ 85分", "强烈推荐", "优先安排面试, 1个工作日内联系"),
    (ACCENT, "70-84分", "推荐面试", "纳入面试名单, 3个工作日内联系"),
    (ACCENT4, "60-69分", "待定", "根据岗位紧急程度决定是否面试"),
    (ACCENT2, "< 60分", "不推荐", "发送感谢信, 纳入人才库"),
]

for i, (color, score, label, action) in enumerate(decisions):
    y = Inches(2.3 + i * 1.1)

    # score badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(2.0), Inches(0.85))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = score
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # arrow
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), y + Inches(0.15), Inches(0.8), Inches(0.5))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()

    # action
    txBox = slide.shapes.add_textbox(Inches(4.2), y + Inches(0.05), Inches(7.5), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = action
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

# stamp area
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(12), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "最终决策: □ 录用  □ 进入下一轮  □ 淘汰  □ 放入人才库"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = ACCENT4
p2 = tf.add_paragraph()
p2.text = "面试官签字: _______________      HR签字: _______________      日期: _______________"
p2.font.size = Pt(12)
p2.font.color.rgb = MID_GRAY
p2.space_before = Pt(12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 ── 尾页
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_bottom_bar(slide)
add_top_accent_line(slide)

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "感谢使用简历筛选模板"
p.font.size = Pt(42)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8))
tf = txBox2.text_frame
p = tf.paragraphs[0]
p.text = "科学筛选 · 客观评估 · 高效决策"
p.font.size = Pt(20)
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# ── Save ───────────────────────────────────────────────────────────
output_path = "C:/Users/SGwj/Desktop/gg/简历筛选模板.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
